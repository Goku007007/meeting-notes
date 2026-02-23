from __future__ import annotations

from email import policy
from email.parser import BytesParser
from pathlib import Path


class UnsupportedFormatError(ValueError):
    pass


class NoExtractableTextError(ValueError):
    pass


SUPPORTED_MIME_BY_EXTENSION: dict[str, set[str]] = {
    ".pdf": {"application/pdf"},
    ".docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
    ".pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    ".xlsx": {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
    ".html": {"text/html"},
    ".htm": {"text/html"},
    ".eml": {"message/rfc822"},
    ".txt": {"text/plain"},
    ".md": {"text/markdown", "text/plain"},
    ".png": {"image/png"},
    ".jpg": {"image/jpeg"},
    ".jpeg": {"image/jpeg"},
    ".webp": {"image/webp"},
}

SUPPORTED_EXTENSIONS = frozenset(SUPPORTED_MIME_BY_EXTENSION.keys())
SUPPORTED_MIME_TYPES = frozenset(
    mime for mime_set in SUPPORTED_MIME_BY_EXTENSION.values() for mime in mime_set
)
GENERIC_MIME_TYPES = {"application/octet-stream", ""}


def validate_supported_upload(filename: str | None, mime_type: str | None) -> tuple[str, str]:
    """
    Validate upload type using both extension and MIME information.
    Returns normalized (extension, mime_type).
    Raises UnsupportedFormatError for unsupported input.
    """
    extension = Path(filename or "").suffix.lower()
    normalized_mime = (mime_type or "").strip().lower()

    if extension not in SUPPORTED_EXTENSIONS:
        raise UnsupportedFormatError(
            f"unsupported file extension '{extension or 'none'}'. "
            "Supported: PDF, DOCX, PPTX, XLSX, HTML, EML, TXT/MD, PNG/JPG/WEBP."
        )

    if normalized_mime and normalized_mime not in SUPPORTED_MIME_TYPES and normalized_mime not in GENERIC_MIME_TYPES:
        raise UnsupportedFormatError(
            f"unsupported media type '{normalized_mime}'. "
            "Supported: PDF, DOCX, PPTX, XLSX, HTML, EML, TXT/MD, PNG/JPG/WEBP."
        )

    expected_mimes = SUPPORTED_MIME_BY_EXTENSION[extension]
    if normalized_mime and normalized_mime not in expected_mimes and normalized_mime not in GENERIC_MIME_TYPES:
        raise UnsupportedFormatError(
            f"file type mismatch: extension '{extension}' is not compatible with mime '{normalized_mime}'"
        )

    resolved_mime = normalized_mime or next(iter(expected_mimes))
    return extension, resolved_mime


def _normalize_text(text: str) -> str:
    return text.replace("\x00", "").strip()


def _read_text_with_fallback(path: Path) -> str:
    for encoding in ("utf-8", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return path.read_text(encoding="utf-8", errors="ignore")


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    parts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            parts.append(page_text)

    text = _normalize_text("\n".join(parts))
    if not text:
        raise NoExtractableTextError("No extractable text found (scanned PDF or image).")
    return text


def _extract_docx(path: Path) -> str:
    from docx import Document as DocxDocument

    doc = DocxDocument(str(path))
    parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    text = _normalize_text("\n".join(parts))
    if not text:
        raise NoExtractableTextError("No extractable text found in DOCX.")
    return text


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text)

    merged = _normalize_text("\n".join(parts))
    if not merged:
        raise NoExtractableTextError("No extractable text found in PPTX.")
    return merged


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(filename=str(path), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in wb.worksheets:
        rows.append(f"[Sheet] {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                rows.append(" | ".join(values))

    text = _normalize_text("\n".join(rows))
    if not text:
        raise NoExtractableTextError("No extractable text found in XLSX.")
    return text


def _extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    html = _read_text_with_fallback(path)
    soup = BeautifulSoup(html, "html.parser")
    text = _normalize_text(soup.get_text(separator="\n"))
    if not text:
        raise NoExtractableTextError("No extractable text found in HTML.")
    return text


def _extract_eml(path: Path) -> str:
    from bs4 import BeautifulSoup

    raw_bytes = path.read_bytes()
    message = BytesParser(policy=policy.default).parsebytes(raw_bytes)

    lines: list[str] = []
    subject = message.get("subject")
    from_header = message.get("from")
    to_header = message.get("to")
    if subject:
        lines.append(f"Subject: {subject}")
    if from_header:
        lines.append(f"From: {from_header}")
    if to_header:
        lines.append(f"To: {to_header}")

    plain_parts: list[str] = []
    html_parts: list[str] = []
    if message.is_multipart():
        for part in message.walk():
            if part.get_content_disposition() == "attachment":
                continue
            ctype = part.get_content_type()
            payload = part.get_content()
            if not isinstance(payload, str):
                continue
            if ctype == "text/plain":
                plain_parts.append(payload)
            elif ctype == "text/html":
                html_parts.append(payload)
    else:
        payload = message.get_content()
        if isinstance(payload, str):
            if message.get_content_type() == "text/html":
                html_parts.append(payload)
            else:
                plain_parts.append(payload)

    body = "\n".join(plain_parts).strip()
    if not body and html_parts:
        soup = BeautifulSoup("\n".join(html_parts), "html.parser")
        body = soup.get_text(separator="\n")
    if body.strip():
        lines.append(body)

    text = _normalize_text("\n\n".join(lines))
    if not text:
        raise NoExtractableTextError("No extractable text found in EML.")
    return text


def _extract_plain_text(path: Path) -> str:
    text = _normalize_text(_read_text_with_fallback(path))
    if not text:
        raise NoExtractableTextError("No extractable text found in text file.")
    return text


def extract_text_from_file(storage_path: str, mime_type: str | None = None) -> str:
    """
    Extract normalized text from supported files.
    Supported: PDF, DOCX, PPTX, XLSX, HTML, EML, TXT/MD.
    Images are accepted at upload but fail extraction unless OCR is added.
    """
    file_path = Path(storage_path)
    if not file_path.exists() or not file_path.is_file():
        raise FileNotFoundError(f"uploaded file not found at {storage_path}")

    extension, _resolved_mime = validate_supported_upload(file_path.name, mime_type)

    if extension in {".txt", ".md"}:
        return _extract_plain_text(file_path)
    if extension == ".pdf":
        return _extract_pdf(file_path)
    if extension == ".docx":
        return _extract_docx(file_path)
    if extension == ".pptx":
        return _extract_pptx(file_path)
    if extension == ".xlsx":
        return _extract_xlsx(file_path)
    if extension in {".html", ".htm"}:
        return _extract_html(file_path)
    if extension == ".eml":
        return _extract_eml(file_path)
    if extension in {".png", ".jpg", ".jpeg", ".webp"}:
        raise NoExtractableTextError("No extractable text found (image uploaded, OCR is not enabled).")

    raise UnsupportedFormatError(
        f"unsupported format for extraction (mime={mime_type or 'unknown'}, ext={extension or 'none'})"
    )

